"""PowerPoint (.pptx) content extractor using markitdown + python-pptx for tables."""

from __future__ import annotations

from pathlib import Path

import pandas as pd
from markitdown import MarkItDown

from extractors.base import BaseExtractor, ExtractedData, ExtractionConfig


class PptxExtractor(BaseExtractor):
    def can_handle(self, file_path: Path) -> bool:
        return file_path.suffix.lower() == ".pptx"

    def extract(self, file_path: Path, config: ExtractionConfig | None = None) -> ExtractedData:
        from pptx import Presentation

        config = config or ExtractionConfig()

        # Use markitdown for rich markdown text extraction
        md = MarkItDown(enable_plugins=False)
        conversion = md.convert(str(file_path))
        full_text = conversion.text_content or ""

        result = ExtractedData(source_path=file_path, text_content=full_text)
        result.dataframes["default"] = pd.DataFrame({"content": [full_text]})

        # Use python-pptx for structured table extraction
        prs = Presentation(str(file_path))
        slide_indices = _parse_slide_range(config.slide_range, len(prs.slides))
        table_counter = 0

        for i, slide in enumerate(prs.slides):
            if i not in slide_indices:
                continue

            for shape in slide.shapes:
                if shape.has_table:
                    tbl = shape.table
                    headers = [cell.text.strip() for cell in tbl.rows[0].cells]
                    rows = [
                        [cell.text.strip() for cell in row.cells]
                        for j, row in enumerate(tbl.rows)
                        if j > 0
                    ]
                    result.dataframes[f"slide_{i + 1}_table_{table_counter}"] = pd.DataFrame(
                        rows, columns=headers
                    )
                    table_counter += 1

        # Extract notes if configured
        if config.include_notes:
            notes_text = []
            for i, slide in enumerate(prs.slides):
                if i in slide_indices and slide.has_notes_slide:
                    notes_text.append(
                        {"slide": i + 1, "notes": slide.notes_slide.notes_text_frame.text}
                    )
            if notes_text:
                result.dataframes["notes"] = pd.DataFrame(notes_text)

        result.metadata["slide_count"] = len(prs.slides)
        return result


def _parse_slide_range(slide_range: str | None, total_slides: int) -> set[int]:
    """Parse slide range string like '1-3,5' into a set of 0-based indices."""
    if not slide_range:
        return set(range(total_slides))

    indices = set()
    for part in slide_range.split(","):
        part = part.strip()
        if "-" in part:
            start, end = part.split("-", 1)
            start_idx = max(0, int(start) - 1)
            end_idx = min(total_slides, int(end))
            indices.update(range(start_idx, end_idx))
        else:
            idx = int(part) - 1
            if 0 <= idx < total_slides:
                indices.add(idx)
    return indices
