"""Script to programmatically create .docx test fixtures.

Run: conda run -n docforge python tests/fixtures/create_fixtures.py
"""

from pathlib import Path

import yaml
from docx import Document
from docx.shared import RGBColor
from pptx import Presentation
from pptx.util import Inches, Pt

FIXTURES_DIR = Path(__file__).parent
TEMPLATES_DIR = FIXTURES_DIR / "templates"
DATA_DIR = FIXTURES_DIR / "data"

TEMPLATES_DIR.mkdir(parents=True, exist_ok=True)
DATA_DIR.mkdir(parents=True, exist_ok=True)


def _add_red_text(paragraph, text):
    """Add a red-colored run to a paragraph."""
    run = paragraph.add_run(text)
    run.font.color.rgb = RGBColor(0xFF, 0x00, 0x00)
    return run


def create_simple_placeholder():
    """Create a doc with one heading and one red text placeholder."""
    doc = Document()
    doc.add_heading("Introduction", level=1)
    para = doc.add_paragraph("The project name is ")
    _add_red_text(para, "Project Name")
    para.add_run(".")
    doc.save(str(TEMPLATES_DIR / "simple_placeholder.docx"))
    print("Created simple_placeholder.docx")


def create_simple_table():
    """Create a doc with one heading and one skeleton table."""
    doc = Document()
    doc.add_heading("Metrics", level=1)
    doc.add_paragraph("The following table shows key metrics:")
    table = doc.add_table(rows=1, cols=3)
    table.style = "Table Grid"
    headers = table.rows[0].cells
    headers[0].text = "Metric"
    headers[1].text = "Value"
    headers[2].text = "Target"
    doc.save(str(TEMPLATES_DIR / "simple_table.docx"))
    print("Created simple_table.docx")


def create_mixed_markers():
    """Create a doc with placeholders, a table, and an LLM prompt marker."""
    doc = Document()

    # Section 1 with placeholders
    doc.add_heading("Project Overview", level=1)
    p1 = doc.add_paragraph("Report prepared by ")
    _add_red_text(p1, "Author")
    p1.add_run(" on ")
    _add_red_text(p1, "Report Date")
    p1.add_run(".")

    # Section 2 with a skeleton table
    doc.add_heading("Key Metrics", level=1)
    doc.add_paragraph("Below are the quarterly metrics:")
    table = doc.add_table(rows=1, cols=3)
    table.style = "Table Grid"
    headers = table.rows[0].cells
    headers[0].text = "Quarter"
    headers[1].text = "Revenue"
    headers[2].text = "Growth"

    # Section 3 with an LLM prompt
    doc.add_heading("Executive Summary", level=1)
    p2 = doc.add_paragraph()
    _add_red_text(
        p2,
        "Provide an executive summary covering the key findings from the quarterly data",
    )

    doc.save(str(TEMPLATES_DIR / "mixed_markers.docx"))
    print("Created mixed_markers.docx")


def create_table_with_sample_data():
    """Create a doc with a skeleton table containing red sample data rows."""
    doc = Document()
    doc.add_heading("Project Status", level=1)
    table = doc.add_table(rows=3, cols=3)
    table.style = "Table Grid"
    # Header row
    headers = table.rows[0].cells
    headers[0].text = "Project"
    headers[1].text = "Status"
    headers[2].text = "Progress"
    # Sample data rows in red
    for row_idx in range(1, 3):
        for col_idx in range(3):
            cell = table.rows[row_idx].cells[col_idx]
            cell.text = ""  # Clear default
            para = cell.paragraphs[0]
            sample_texts = [
                ["Alpha", "Active", "75%"],
                ["Beta", "Planning", "20%"],
            ]
            _add_red_text(para, sample_texts[row_idx - 1][col_idx])

    doc.save(str(TEMPLATES_DIR / "table_with_sample_data.docx"))
    print("Created table_with_sample_data.docx")


def create_data_fixtures():
    """Create test data files."""
    import json

    import pandas as pd

    # Excel with two sheets
    with pd.ExcelWriter(str(DATA_DIR / "metrics.xlsx"), engine="openpyxl") as writer:
        df1 = pd.DataFrame(
            {
                "Quarter": ["Q1", "Q2", "Q3", "Q4"],
                "Revenue": [100000, 120000, 115000, 140000],
                "Growth": [0.05, 0.20, -0.04, 0.22],
            }
        )
        df1.to_excel(writer, sheet_name="Revenue", index=False)
        df2 = pd.DataFrame(
            {
                "Metric": ["Users", "Sessions", "Conversions"],
                "Value": [5000, 25000, 500],
                "Target": [4500, 20000, 450],
            }
        )
        df2.to_excel(writer, sheet_name="KPIs", index=False)
    print("Created metrics.xlsx")

    # CSV
    df_csv = pd.DataFrame(
        {
            "Project": ["Alpha", "Beta", "Gamma", "Delta"],
            "Status": ["Active", "Planning", "Complete", "Active"],
            "Progress": [75, 20, 100, 45],
            "Lead": ["Alice", "Bob", "Carol", "Dave"],
        }
    )
    df_csv.to_csv(str(DATA_DIR / "project_status.csv"), index=False)
    print("Created project_status.csv")

    # JSON
    config = {
        "project": {
            "name": "DocForge Demo",
            "date": "2026-03-03",
            "version": "1.0",
        },
        "settings": {"author": "Test Author", "department": "Engineering"},
    }
    with open(DATA_DIR / "config.json", "w") as f:
        json.dump(config, f, indent=2)
    print("Created config.json")


def _create_raw_pdf(page_texts: list[str]) -> bytes:
    """Create a minimal valid PDF with the given page texts (no external dependencies)."""
    objects = []
    obj_num = 1

    # Catalog
    catalog_num = obj_num
    objects.append(f"{obj_num} 0 obj\n<< /Type /Catalog /Pages {obj_num + 1} 0 R >>\nendobj")
    obj_num += 1

    # Pages (parent)
    pages_num = obj_num
    kid_nums = []
    obj_num += 1  # Reserve slot, fill later

    # Font
    font_num = obj_num
    objects.append(
        f"{obj_num} 0 obj\n<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>\nendobj"
    )
    obj_num += 1

    # Create each page + contents
    for text in page_texts:
        page_num = obj_num
        kid_nums.append(page_num)
        content_num = obj_num + 1

        stream = f"BT /F1 11 Tf 72 720 Td ({text}) Tj ET"
        objects.append(
            f"{page_num} 0 obj\n"
            f"<< /Type /Page /Parent {pages_num} 0 R /MediaBox [0 0 612 792]\n"
            f"   /Contents {content_num} 0 R "
            f"/Resources << /Font << /F1 {font_num} 0 R >> >> >>\n"
            f"endobj"
        )
        objects.append(
            f"{content_num} 0 obj\n<< /Length {len(stream)} >>\nstream\n{stream}\nendstream\nendobj"
        )
        obj_num += 2

    # Now insert the Pages object at position 1
    kids_str = " ".join(f"{n} 0 R" for n in kid_nums)
    pages_obj = (
        f"{pages_num} 0 obj\n<< /Type /Pages /Kids [{kids_str}] /Count {len(page_texts)} >>\nendobj"
    )
    objects.insert(1, pages_obj)

    # Build PDF
    body = "\n".join(objects)
    xref_offset = len(b"%PDF-1.4\n") + len(body.encode()) + 1
    trailer = f"<< /Root {catalog_num} 0 R /Size {obj_num} >>"
    pdf = (
        f"%PDF-1.4\n{body}\nxref\n0 1\n0000000000 65535 f \n"
        f"trailer\n{trailer}\nstartxref\n{xref_offset}\n%%EOF"
    )
    return pdf.encode()


def create_phase2_data_fixtures():
    """Create Phase 2 data fixtures: .txt, .yaml, .docx, .pptx, .pdf."""

    # --- data/sample.txt ---
    sample_text = (
        "Quarterly Performance Report\n"
        "\n"
        "The first quarter showed strong momentum across all business units. "
        "Revenue exceeded targets by 12%, driven primarily by new customer "
        "acquisitions in the enterprise segment. Operating margins improved "
        "to 18.5%, up from 16.2% in the prior quarter.\n"
        "\n"
        "The product team shipped three major releases during the quarter, "
        "including the long-awaited analytics dashboard and a revamped "
        "onboarding flow. Customer satisfaction scores rose to 4.6 out of 5, "
        "reflecting the positive reception of these updates.\n"
        "\n"
        "Looking ahead, the focus for Q2 will be on international expansion "
        "and deepening integrations with key partners. The hiring pipeline "
        "remains healthy, with 15 new positions approved for engineering "
        "and customer success teams.\n"
    )
    (DATA_DIR / "sample.txt").write_text(sample_text, encoding="utf-8")
    print("Created data/sample.txt")

    # --- data/config.yaml ---
    config_yaml = {
        "project": {
            "name": "DocForge Demo",
            "date": "2026-03-03",
        },
        "settings": {
            "author": "Test Author",
            "theme": "dark",
        },
    }
    with open(DATA_DIR / "config.yaml", "w", encoding="utf-8") as f:
        yaml.dump(config_yaml, f, default_flow_style=False, sort_keys=False)
    print("Created data/config.yaml")

    # --- data/source_doc.docx ---
    doc = Document()
    doc.add_paragraph(
        "This document serves as a supplementary data source for document "
        "generation. It contains contextual information about the project "
        "team and their responsibilities."
    )
    doc.add_paragraph(
        "The team has been restructured to align with the new strategic "
        "priorities. Each department now has a dedicated liaison responsible "
        "for cross-functional coordination."
    )
    doc.add_paragraph(
        "The following table summarises the current team assignments and departmental affiliations."
    )

    table = doc.add_table(rows=3, cols=3)
    table.style = "Table Grid"
    # Header row
    for idx, header in enumerate(["Name", "Role", "Department"]):
        table.rows[0].cells[idx].text = header
    # Data rows
    data_rows = [
        ("Alice Johnson", "Project Lead", "Engineering"),
        ("Bob Smith", "Data Analyst", "Analytics"),
    ]
    for row_idx, row_data in enumerate(data_rows, start=1):
        for col_idx, value in enumerate(row_data):
            table.rows[row_idx].cells[col_idx].text = value

    doc.save(str(DATA_DIR / "source_doc.docx"))
    print("Created data/source_doc.docx")

    # --- data/slides.pptx ---
    prs = Presentation()

    # Slide 1: title + body text
    slide_layout_title = prs.slide_layouts[1]  # Title and Content
    slide1 = prs.slides.add_slide(slide_layout_title)
    slide1.shapes.title.text = "Overview"
    slide1.placeholders[1].text = (
        "This presentation provides a high-level overview of the project "
        "status and key metrics for the current reporting period."
    )

    # Slide 2: title + table
    slide_layout_blank = prs.slide_layouts[5]  # Blank
    slide2 = prs.slides.add_slide(slide_layout_blank)

    # Add title text box
    tx_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
    tf = tx_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Data"
    p.font.size = Pt(28)
    p.font.bold = True

    # Add table: header + 2 data rows = 3 rows, 2 cols
    table_shape = slide2.shapes.add_table(
        rows=3,
        cols=2,
        left=Inches(1.5),
        top=Inches(1.8),
        width=Inches(6),
        height=Inches(2),
    )
    tbl = table_shape.table
    tbl.cell(0, 0).text = "Metric"
    tbl.cell(0, 1).text = "Value"
    tbl.cell(1, 0).text = "Revenue"
    tbl.cell(1, 1).text = "$140,000"
    tbl.cell(2, 0).text = "Growth"
    tbl.cell(2, 1).text = "22%"

    prs.save(str(DATA_DIR / "slides.pptx"))
    print("Created data/slides.pptx")

    # --- data/sample.pdf ---
    # Create a minimal 2-page PDF using raw PDF syntax (no pymupdf dependency)
    page1_text = (
        "Quarterly Performance Summary - Page 1  "
        "The organisation achieved record revenue in Q4, surpassing the "
        "annual target by 8 percent."
    )
    page2_text = (
        "Quarterly Performance Summary - Page 2  "
        "Marketing efforts in the quarter focused on content-led growth."
    )
    pdf_content = _create_raw_pdf([page1_text, page2_text])
    (DATA_DIR / "sample.pdf").write_bytes(pdf_content)
    print("Created data/sample.pdf")


def create_conditional_template():
    """Create a template for testing conditional section removal.

    Structure:
      - Heading 1: "Report"
      - Paragraph with red placeholder: "Title"
      - Heading 1: "Optional Section"
      - Paragraph with red LLM prompt: "Summarize the optional data if available"
      - Heading 1: "Summary"
      - Regular text paragraph
    """
    doc = Document()

    doc.add_heading("Report", level=1)
    p1 = doc.add_paragraph()
    _add_red_text(p1, "Title")

    doc.add_heading("Optional Section", level=1)
    p2 = doc.add_paragraph()
    _add_red_text(p2, "Summarize the optional data if available")

    doc.add_heading("Summary", level=1)
    doc.add_paragraph("This section contains the final summary and conclusions of the report.")

    doc.save(str(TEMPLATES_DIR / "conditional_template.docx"))
    print("Created templates/conditional_template.docx")


if __name__ == "__main__":
    create_simple_placeholder()
    create_simple_table()
    create_mixed_markers()
    create_table_with_sample_data()
    create_data_fixtures()
    create_phase2_data_fixtures()
    create_conditional_template()
    print("\nAll fixtures created successfully!")
