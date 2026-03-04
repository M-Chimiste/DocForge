"""Tests for Phase 2 data extractors: text, yaml, docx, pptx, pdf."""

from __future__ import annotations

from pathlib import Path

import pytest

from extractors.base import ExtractionConfig
from extractors.docx_extractor import DocxContentExtractor
from extractors.pdf_extractor import PdfExtractor
from extractors.pptx_extractor import PptxExtractor
from extractors.text_extractor import TextExtractor
from extractors.yaml_extractor import YamlExtractor

# ---------------------------------------------------------------------------
# Helpers to create temp fixture files
# ---------------------------------------------------------------------------


def _create_text_file(tmp_path: Path, name: str = "sample.txt", content: str | None = None) -> Path:
    text = content or "Line one\nLine two\nLine three"
    path = tmp_path / name
    path.write_text(text, encoding="utf-8")
    return path


def _create_yaml_file(tmp_path: Path, name: str = "data.yaml", content: str | None = None) -> Path:
    yaml_text = content or ("name: DocForge\nversion: 1.0\nactive: true\n")
    path = tmp_path / name
    path.write_text(yaml_text, encoding="utf-8")
    return path


def _create_nested_yaml_file(tmp_path: Path) -> Path:
    yaml_text = (
        "project:\n"
        "  name: DocForge\n"
        "  version: 2.0\n"
        "  team:\n"
        "    - name: Alice\n"
        "      role: Lead\n"
        "    - name: Bob\n"
        "      role: Dev\n"
    )
    path = tmp_path / "nested.yaml"
    path.write_text(yaml_text, encoding="utf-8")
    return path


def _create_docx_file(tmp_path: Path) -> Path:
    from docx import Document

    doc = Document()
    doc.add_heading("Test Document", level=1)
    doc.add_paragraph("First paragraph of content.")
    doc.add_paragraph("Second paragraph of content.")

    table = doc.add_table(rows=3, cols=2)
    table.style = "Table Grid"
    table.rows[0].cells[0].text = "Name"
    table.rows[0].cells[1].text = "Score"
    table.rows[1].cells[0].text = "Alice"
    table.rows[1].cells[1].text = "95"
    table.rows[2].cells[0].text = "Bob"
    table.rows[2].cells[1].text = "87"

    path = tmp_path / "test_doc.docx"
    doc.save(str(path))
    return path


def _create_pptx_file(tmp_path: Path) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()

    # Slide 1 -- title + body text
    slide1_layout = prs.slide_layouts[1]
    slide1 = prs.slides.add_slide(slide1_layout)
    slide1.shapes.title.text = "Slide One Title"
    slide1.placeholders[1].text = "Slide one body text."

    # Slide 2 -- text + table
    slide2_layout = prs.slide_layouts[5]
    slide2 = prs.slides.add_slide(slide2_layout)

    tx_box = slide2.shapes.add_textbox(Inches(1), Inches(0.5), Inches(5), Inches(1))
    tx_box.text_frame.text = "Slide two text."

    rows, cols = 3, 2
    tbl_shape = slide2.shapes.add_table(rows, cols, Inches(1), Inches(2), Inches(5), Inches(2))
    tbl = tbl_shape.table
    tbl.rows[0].cells[0].text = "Item"
    tbl.rows[0].cells[1].text = "Count"
    tbl.rows[1].cells[0].text = "Widgets"
    tbl.rows[1].cells[1].text = "42"
    tbl.rows[2].cells[0].text = "Gizmos"
    tbl.rows[2].cells[1].text = "17"

    path = tmp_path / "test_pres.pptx"
    prs.save(str(path))
    return path


def _create_pdf_file(tmp_path: Path) -> Path:
    """Create a minimal 2-page PDF using reportlab or raw PDF bytes."""
    # Create a minimal valid PDF with two pages using raw PDF syntax
    pdf_content = b"""%PDF-1.4
1 0 obj
<< /Type /Catalog /Pages 2 0 R >>
endobj
2 0 obj
<< /Type /Pages /Kids [3 0 R 6 0 R] /Count 2 >>
endobj
3 0 obj
<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]
   /Contents 4 0 R /Resources << /Font << /F1 5 0 R >> >> >>
endobj
4 0 obj
<< /Length 44 >>
stream
BT /F1 12 Tf 72 720 Td (Page one content.) Tj ET
endstream
endobj
5 0 obj
<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica >>
endobj
6 0 obj
<< /Type /Page /Parent 2 0 R /MediaBox [0 0 612 792]
   /Contents 7 0 R /Resources << /Font << /F1 5 0 R >> >> >>
endobj
7 0 obj
<< /Length 44 >>
stream
BT /F1 12 Tf 72 720 Td (Page two content.) Tj ET
endstream
endobj
xref
0 8
0000000000 65535 f
0000000009 00000 n
0000000058 00000 n
0000000115 00000 n
0000000266 00000 n
0000000360 00000 n
0000000441 00000 n
0000000592 00000 n
trailer
<< /Root 1 0 R /Size 8 >>
startxref
686
%%EOF"""
    path = tmp_path / "test_doc.pdf"
    path.write_bytes(pdf_content)
    return path


# ---------------------------------------------------------------------------
# TestTextExtractor
# ---------------------------------------------------------------------------


class TestTextExtractor:
    def setup_method(self):
        self.extractor = TextExtractor()

    def test_can_handle_txt(self):
        assert self.extractor.can_handle(Path("readme.txt")) is True

    def test_can_handle_md(self):
        assert self.extractor.can_handle(Path("notes.md")) is True

    def test_cannot_handle_json(self):
        assert self.extractor.can_handle(Path("data.json")) is False

    def test_cannot_handle_csv(self):
        assert self.extractor.can_handle(Path("data.csv")) is False

    def test_extract_reads_text(self, tmp_path):
        content = "Hello World\nSecond line"
        path = _create_text_file(tmp_path, content=content)
        result = self.extractor.extract(path)
        assert result.text_content == content

    def test_extract_default_dataframe(self, tmp_path):
        content = "Some text"
        path = _create_text_file(tmp_path, content=content)
        result = self.extractor.extract(path)
        assert "default" in result.dataframes
        assert result.dataframes["default"]["content"].iloc[0] == content

    def test_metadata_char_count(self, tmp_path):
        content = "abc"
        path = _create_text_file(tmp_path, content=content)
        result = self.extractor.extract(path)
        assert result.metadata["char_count"] == 3

    def test_metadata_line_count(self, tmp_path):
        content = "Line 1\nLine 2\nLine 3"
        path = _create_text_file(tmp_path, content=content)
        result = self.extractor.extract(path)
        assert result.metadata["line_count"] == 3

    def test_text_content_is_set(self, tmp_path):
        content = "Non-empty text"
        path = _create_text_file(tmp_path, content=content)
        result = self.extractor.extract(path)
        assert result.text_content is not None
        assert len(result.text_content) > 0

    def test_source_path_set(self, tmp_path):
        path = _create_text_file(tmp_path)
        result = self.extractor.extract(path)
        assert result.source_path == path


# ---------------------------------------------------------------------------
# TestYamlExtractor
# ---------------------------------------------------------------------------


class TestYamlExtractor:
    def setup_method(self):
        self.extractor = YamlExtractor()

    def test_can_handle_yaml(self):
        assert self.extractor.can_handle(Path("config.yaml")) is True

    def test_can_handle_yml(self):
        assert self.extractor.can_handle(Path("config.yml")) is True

    def test_cannot_handle_json(self):
        assert self.extractor.can_handle(Path("data.json")) is False

    def test_cannot_handle_txt(self):
        assert self.extractor.can_handle(Path("notes.txt")) is False

    def test_extract_flat_yaml(self, tmp_path):
        path = _create_yaml_file(tmp_path)
        result = self.extractor.extract(path)
        df = result.dataframes["default"]
        assert "name" in df.columns
        assert df["name"].iloc[0] == "DocForge"
        assert df["version"].iloc[0] == 1.0

    def test_extract_nested_with_yaml_path(self, tmp_path):
        path = _create_nested_yaml_file(tmp_path)
        config = ExtractionConfig(yaml_path="project.team")
        result = self.extractor.extract(path, config=config)
        df = result.dataframes["project.team"]
        assert len(df) == 2
        assert df["name"].iloc[0] == "Alice"
        assert df["role"].iloc[1] == "Dev"

    def test_metadata_has_raw_key(self, tmp_path):
        path = _create_yaml_file(tmp_path)
        result = self.extractor.extract(path)
        assert "raw" in result.metadata
        assert result.metadata["raw"]["name"] == "DocForge"

    def test_extract_nested_yaml_path_single_value(self, tmp_path):
        path = _create_nested_yaml_file(tmp_path)
        config = ExtractionConfig(yaml_path="project.name")
        result = self.extractor.extract(path, config=config)
        df = result.dataframes["project.name"]
        assert df["value"].iloc[0] == "DocForge"

    def test_extract_yaml_path_invalid_key_raises(self, tmp_path):
        path = _create_yaml_file(tmp_path)
        config = ExtractionConfig(yaml_path="nonexistent.key")
        with pytest.raises(KeyError):
            self.extractor.extract(path, config=config)


# ---------------------------------------------------------------------------
# TestDocxContentExtractor
# ---------------------------------------------------------------------------


class TestDocxContentExtractor:
    def setup_method(self):
        self.extractor = DocxContentExtractor()

    def test_can_handle_docx(self):
        assert self.extractor.can_handle(Path("report.docx")) is True

    def test_cannot_handle_pdf(self):
        assert self.extractor.can_handle(Path("report.pdf")) is False

    def test_cannot_handle_pptx(self):
        assert self.extractor.can_handle(Path("slides.pptx")) is False

    def test_extract_body_text(self, tmp_path):
        path = _create_docx_file(tmp_path)
        result = self.extractor.extract(path)
        # markitdown produces markdown — check key content is present
        assert result.text_content is not None
        assert "First paragraph" in result.text_content
        assert "Second paragraph" in result.text_content

    def test_extract_tables_as_dataframes(self, tmp_path):
        path = _create_docx_file(tmp_path)
        result = self.extractor.extract(path)
        assert "table_0" in result.dataframes
        table_df = result.dataframes["table_0"]
        assert list(table_df.columns) == ["Name", "Score"]
        assert len(table_df) == 2
        assert table_df["Name"].iloc[0] == "Alice"
        assert table_df["Score"].iloc[1] == "87"

    def test_text_content_is_set(self, tmp_path):
        path = _create_docx_file(tmp_path)
        result = self.extractor.extract(path)
        assert result.text_content is not None
        assert len(result.text_content) > 0

    def test_default_dataframe_has_body(self, tmp_path):
        path = _create_docx_file(tmp_path)
        result = self.extractor.extract(path)
        assert "default" in result.dataframes
        body = result.dataframes["default"]["content"].iloc[0]
        assert "First paragraph" in body

    def test_metadata_counts(self, tmp_path):
        path = _create_docx_file(tmp_path)
        result = self.extractor.extract(path)
        assert result.metadata["paragraph_count"] > 0
        assert result.metadata["table_count"] == 1


# ---------------------------------------------------------------------------
# TestPptxExtractor
# ---------------------------------------------------------------------------


class TestPptxExtractor:
    def setup_method(self):
        self.extractor = PptxExtractor()

    def test_can_handle_pptx(self):
        assert self.extractor.can_handle(Path("deck.pptx")) is True

    def test_cannot_handle_docx(self):
        assert self.extractor.can_handle(Path("doc.docx")) is False

    def test_cannot_handle_pdf(self):
        assert self.extractor.can_handle(Path("file.pdf")) is False

    def test_extract_text(self, tmp_path):
        path = _create_pptx_file(tmp_path)
        result = self.extractor.extract(path)
        assert result.text_content is not None
        # markitdown produces markdown — check key content is present
        assert "Slide One Title" in result.text_content

    def test_extract_tables(self, tmp_path):
        path = _create_pptx_file(tmp_path)
        result = self.extractor.extract(path)
        table_keys = [k for k in result.dataframes if k.startswith("slide_")]
        assert len(table_keys) >= 1
        table_df = result.dataframes[table_keys[0]]
        assert "Item" in table_df.columns
        assert "Count" in table_df.columns
        assert len(table_df) == 2

    def test_metadata_has_slide_count(self, tmp_path):
        path = _create_pptx_file(tmp_path)
        result = self.extractor.extract(path)
        assert result.metadata["slide_count"] == 2

    def test_default_dataframe_exists(self, tmp_path):
        path = _create_pptx_file(tmp_path)
        result = self.extractor.extract(path)
        assert "default" in result.dataframes


# ---------------------------------------------------------------------------
# TestPdfExtractor
# ---------------------------------------------------------------------------


class TestPdfExtractor:
    def setup_method(self):
        self.extractor = PdfExtractor()

    def test_can_handle_pdf(self):
        assert self.extractor.can_handle(Path("document.pdf")) is True

    def test_cannot_handle_docx(self):
        assert self.extractor.can_handle(Path("document.docx")) is False

    def test_cannot_handle_txt(self):
        assert self.extractor.can_handle(Path("notes.txt")) is False

    def test_extract_text(self, tmp_path):
        path = _create_pdf_file(tmp_path)
        result = self.extractor.extract(path)
        assert result.text_content is not None
        assert "Page one content" in result.text_content
        assert "Page two content" in result.text_content

    def test_default_dataframe_has_full_text(self, tmp_path):
        path = _create_pdf_file(tmp_path)
        result = self.extractor.extract(path)
        assert "default" in result.dataframes
        full_text = result.dataframes["default"]["content"].iloc[0]
        assert "Page one content" in full_text

    def test_text_content_is_set(self, tmp_path):
        path = _create_pdf_file(tmp_path)
        result = self.extractor.extract(path)
        assert result.text_content is not None
        assert len(result.text_content) > 0

    def test_source_path_set(self, tmp_path):
        path = _create_pdf_file(tmp_path)
        result = self.extractor.extract(path)
        assert result.source_path == path
